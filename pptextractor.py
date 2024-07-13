from pptx import Presentation

# Load your presentation
prs = Presentation('SlideSpeak Crypto Investing Short Master.pptx')

# Open a file to write the text
with open('SlideSpeak Crypto Investing Short Master.txt', 'w', encoding='utf-8') as f:
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(shape.text + '\n')

print("Text extraction complete. Check presentation_text.txt for the output.")